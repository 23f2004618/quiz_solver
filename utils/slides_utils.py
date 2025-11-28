import base64
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt

def generate_slides(slides_data: dict) -> str:
    """
    Generate a PowerPoint presentation from a dictionary of slide data.
    
    Args:
        slides_data: dict containing 'title' and 'slides' list.
                     Example:
                     {
                         "title": "Presentation Title",
                         "subtitle": "Subtitle here",
                         "slides": [
                             {"title": "Slide 1", "content": "Bullet 1\nBullet 2"},
                             {"title": "Slide 2", "content": "Some text content"}
                         ]
                     }
    
    Returns:
        base64 data URI string (data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,...)
    """
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = slides_data.get("title", "Presentation")
    subtitle.text = slides_data.get("subtitle", "")

    # Content Slides
    bullet_slide_layout = prs.slide_layouts[1]

    for slide_info in slides_data.get("slides", []):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slide_info.get("title", "Slide")
        tf = body_shape.text_frame
        
        content = slide_info.get("content", "")
        if isinstance(content, list):
            # If content is a list, treat as bullets
            for i, item in enumerate(content):
                if i == 0:
                    tf.text = str(item)
                else:
                    p = tf.add_paragraph()
                    p.text = str(item)
        else:
            # If content is string, split by newlines for bullets
            lines = str(content).split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    tf.text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line

    # Save to buffer
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    b64_data = base64.b64encode(buffer.getvalue()).decode()
    return f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_data}"
